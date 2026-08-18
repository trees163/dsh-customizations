# vscode-doc-convert.py — docx/xlsx/pptx → HTML 预览转换（供 dsh-host-files 调用）
import sys

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

import html


def esc(s):
    return html.escape(str(s))


def docx_html(path):
    from docx import Document
    doc = Document(path)
    parts = []
    for p in doc.paragraphs:
        t = p.text.strip()
        if t:
            parts.append("<p style='margin:4px 0'>" + esc(t) + "</p>")
    for tb in doc.tables:
        parts.append("<table border='1' cellspacing='0' cellpadding='4' style='border-collapse:collapse;width:100%;margin:8px 0'>")
        for row in tb.rows:
            cells = "".join("<td style='padding:4px 8px'>" + esc(c.text) + "</td>" for c in row.cells)
            parts.append("<tr>" + cells + "</tr>")
        parts.append("</table>")
    return "".join(parts) or "<p>(空文档)</p>"


def xlsx_html(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append("<h3 style='margin:12px 0 6px'>" + esc(ws.title) + "</h3>")
        parts.append("<table border='1' cellspacing='0' cellpadding='4' style='border-collapse:collapse'>")
        rows = 0
        for row in ws.iter_rows(values_only=True):
            cells = "".join("<td style='padding:2px 8px;white-space:nowrap'>" + esc("" if v is None else v) + "</td>" for v in row)
            parts.append("<tr>" + cells + "</tr>")
            rows += 1
            if rows >= 300:
                parts.append("<tr><td style='padding:4px 8px;color:#8a919c'>…（仅显示前 300 行）</td></tr>")
                break
        parts.append("</table>")
    return "".join(parts) or "<p>(空表格)</p>"


def pptx_html(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append("<h3 style='margin:14px 0 6px'>第 " + str(i) + " 页</h3>")
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append("<p style='margin:4px 0'>" + esc(t) + "</p>")
    return "".join(parts) or "<p>(空演示文稿)</p>"


def main():
    path = sys.argv[1]
    ext = path.rsplit(".", 1)[-1].lower() if "." in path else ""
    try:
        if ext == "docx":
            out = docx_html(path)
        elif ext == "xlsx":
            out = xlsx_html(path)
        elif ext == "pptx":
            out = pptx_html(path)
        else:
            raise ValueError("unsupported extension: " + ext)
        print("<div style='font-family:system-ui;font-size:13px;color:#24292f;line-height:1.6'>" + out + "</div>")
    except Exception as e:
        print("<div style='font-family:system-ui;font-size:13px;color:#f97583'>转换失败: " + esc(str(e)) + "</div>")


if __name__ == "__main__":
    main()
